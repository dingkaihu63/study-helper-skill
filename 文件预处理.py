"""
文档预处理工具 - 将 PDF/Word/TXT 等转换为 AI Agent 结构化输入
支持：PDF、DOCX、TXT、MD、HTML、PPT、XLSX
"""

import os
import re
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Tuple
from dataclasses import dataclass, asdict
from enum import Enum
import tempfile
import warnings

# 核心依赖
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
except ImportError:
    markdown = None


class ContentType(Enum):
    """内容类型枚举"""
    TEXT = "text"
    HEADING = "heading"
    TABLE = "table"
    IMAGE = "image"
    LIST = "list"
    CODE = "code"
    METADATA = "metadata"
    PAGE_BREAK = "page_break"


@dataclass
class Chunk:
    """文档分块数据结构"""
    id: str
    type: str
    content: str
    metadata: Dict[str, Any]
    page_num: Optional[int] = None
    section_title: Optional[str] = None
    word_count: int = 0
    char_count: int = 0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "id": self.id,
            "type": self.type,
            "content": self.content,
            "metadata": self.metadata,
            "page_num": self.page_num,
            "section_title": self.section_title,
            "word_count": self.word_count,
            "char_count": self.char_count
        }


@dataclass
class ProcessedDocument:
    """处理后的文档结构"""
    source_file: str
    file_type: str
    total_pages: int
    total_chunks: int
    metadata: Dict[str, Any]
    chunks: List[Chunk]
    summary: Dict[str, Any]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "source_file": self.source_file,
            "file_type": self.file_type,
            "total_pages": self.total_pages,
            "total_chunks": self.total_chunks,
            "metadata": self.metadata,
            "summary": self.summary,
            "chunks": [c.to_dict() for c in self.chunks]
        }

    def to_json(self, indent: int = 2) -> str:
        return json.dumps(self.to_dict(), ensure_ascii=False, indent=indent)

    def to_agent_prompt(self, max_chunks: Optional[int] = None) -> str:
        """转换为 AI Agent 可直接处理的提示文本"""
        lines = [
            f"# 文档: {self.source_file}",
            f"类型: {self.file_type} | 页数: {self.total_pages} | 块数: {self.total_chunks}",
            "",
            "## 文档元数据",
            json.dumps(self.metadata, ensure_ascii=False, indent=2),
            "",
            "## 文档内容",
        ]

        chunks_to_use = self.chunks[:max_chunks] if max_chunks else self.chunks

        for chunk in chunks_to_use:
            lines.append(f"\n### [{chunk.type.upper()}] {chunk.section_title or '未命名段落'}")
            if chunk.page_num:
                lines.append(f"*页码: {chunk.page_num}*")
            lines.append(f"```")
            lines.append(chunk.content)
            lines.append(f"```")
            lines.append(f"*字数: {chunk.word_count} | 字符: {chunk.char_count}*")

        if max_chunks and len(self.chunks) > max_chunks:
            lines.append(f"\n... (还有 {len(self.chunks) - max_chunks} 个块未显示)")

        return "\n".join(lines)


class TextSplitter:
    """智能文本分块器"""

    def __init__(
            self,
            chunk_size: int = 800,
            chunk_overlap: int = 100,
            separators: Optional[List[str]] = None
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or [
            "\n\n",  # 段落
            "\n",  # 换行
            "。",  # 句号
            "．",  # 全角点
            ". ",  # 英文句号+空格
            "；",  # 分号
            ";",  # 英文分号
            "，",  # 逗号
            ",",  # 英文逗号
            " ",  # 空格
            ""  # 字符
        ]

    def split_text(self, text: str) -> List[str]:
        """递归字符文本分割"""
        if len(text) <= self.chunk_size:
            return [text] if text.strip() else []

        for separator in self.separators:
            if separator == "":
                # 字符级分割
                chunks = []
                for i in range(0, len(text), self.chunk_size - self.chunk_overlap):
                    chunk = text[i:i + self.chunk_size]
                    if chunk.strip():
                        chunks.append(chunk)
                return chunks

            if separator in text:
                parts = text.split(separator)
                chunks = []
                current_chunk = ""

                for part in parts:
                    if not part.strip():
                        continue

                    test_chunk = current_chunk + separator + part if current_chunk else part

                    if len(test_chunk) <= self.chunk_size:
                        current_chunk = test_chunk
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = part

                if current_chunk:
                    chunks.append(current_chunk.strip())

                # 处理重叠
                if len(chunks) > 1 and self.chunk_overlap > 0:
                    overlapped_chunks = []
                    for i, chunk in enumerate(chunks):
                        if i > 0 and len(chunk) > self.chunk_overlap:
                            prev_chunk = chunks[i - 1]
                            overlap_text = prev_chunk[-self.chunk_overlap:] if len(
                                prev_chunk) > self.chunk_overlap else prev_chunk
                            chunk = overlap_text + chunk
                        overlapped_chunks.append(chunk)
                    return overlapped_chunks

                return [c for c in chunks if c.strip()]

        return [text[:self.chunk_size]]


class DocumentProcessor:
    """文档处理器主类"""

    def __init__(
            self,
            chunk_size: int = 800,
            chunk_overlap: int = 100,
            extract_images: bool = False,
            image_output_dir: Optional[str] = None,
            ocr_enabled: bool = False
    ):
        self.splitter = TextSplitter(chunk_size, chunk_overlap)
        self.extract_images = extract_images
        self.image_output_dir = image_output_dir or tempfile.mkdtemp()
        self.ocr_enabled = ocr_enabled
        self._current_section = "文档开头"

    def _generate_id(self, content: str, prefix: str = "") -> str:
        """生成内容唯一ID"""
        hash_obj = hashlib.md5(content.encode('utf-8'))
        return f"{prefix}{hash_obj.hexdigest()[:12]}"

    def _create_chunk(
            self,
            content: str,
            chunk_type: str,
            metadata: Dict[str, Any],
            page_num: Optional[int] = None
    ) -> Chunk:
        """创建 Chunk 对象"""
        return Chunk(
            id=self._generate_id(content, f"{chunk_type}_"),
            type=chunk_type,
            content=content.strip(),
            metadata=metadata,
            page_num=page_num,
            section_title=self._current_section,
            word_count=len(content.split()),
            char_count=len(content)
        )

    def _clean_text(self, text: str) -> str:
        """清理文本"""
        # 移除多余空白
        text = re.sub(r'\s+', ' ', text)
        # 移除控制字符
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        # 修复换行
        text = re.sub(r'([。！？；])\s+', r'\1\n', text)
        return text.strip()

    def _extract_table_text(self, table_data: List[List[str]]) -> str:
        """将表格转换为文本"""
        if not table_data:
            return ""

        lines = []
        for row in table_data:
            lines.append(" | ".join(str(cell) for cell in row if cell))
        return "\n".join(lines)

    def process_pdf(self, file_path: str) -> ProcessedDocument:
        """处理 PDF 文件"""
        if fitz is None:
            raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")

        doc = fitz.open(file_path)
        chunks = []
        all_text = []

        metadata = {
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "creator": doc.metadata.get("creator", ""),
            "page_count": len(doc)
        }

        for page_num, page in enumerate(doc, 1):
            # 提取文本块（保持布局）
            blocks = page.get_text("blocks")

            for block in blocks:
                if len(block) < 7:
                    continue

                x0, y0, x1, y1, text, block_no, block_type = block

                if not text.strip():
                    continue

                text = self._clean_text(text)

                # 检测标题（基于字体大小和位置）
                is_heading = False
                if y1 < 150 and len(text) < 100:  # 页面顶部短文本
                    is_heading = True
                    self._current_section = text

                # 检测表格
                if block_type == 1 or self._looks_like_table(text):
                    chunk_type = ContentType.TABLE.value
                elif is_heading:
                    chunk_type = ContentType.HEADING.value
                else:
                    chunk_type = ContentType.TEXT.value

                # 如果文本太长，进行分块
                if len(text) > self.splitter.chunk_size:
                    sub_chunks = self.splitter.split_text(text)
                    for sub_text in sub_chunks:
                        chunks.append(self._create_chunk(
                            sub_text, chunk_type,
                            {"bbox": [x0, y0, x1, y1], "block_no": block_no},
                            page_num
                        ))
                else:
                    chunks.append(self._create_chunk(
                        text, chunk_type,
                        {"bbox": [x0, y0, x1, y1], "block_no": block_no},
                        page_num
                    ))

                all_text.append(text)

            # 提取图片
            if self.extract_images:
                img_chunks = self._extract_pdf_images(page, page_num, file_path)
                chunks.extend(img_chunks)

        doc.close()

        return self._build_document(file_path, "pdf", len(doc), chunks, metadata, all_text)

    def _looks_like_table(self, text: str) -> bool:
        """启发式判断是否为表格文本"""
        lines = text.strip().split('\n')
        if len(lines) < 2:
            return False

        # 检查是否有大量制表符或竖线分隔
        tab_count = text.count('\t') + text.count(' | ') + text.count('|')
        return tab_count > 3 or (len(lines) > 2 and all(len(l.split()) > 2 for l in lines[:3]))

    def _extract_pdf_images(self, page, page_num: int, file_path: str) -> List[Chunk]:
        """提取 PDF 中的图片"""
        chunks = []
        image_list = page.get_images()

        for img_index, img in enumerate(image_list, 1):
            xref = img[0]
            try:
                base_image = page.parent.extract_image(xref)
                if base_image:
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]

                    # 保存图片
                    img_filename = f"{Path(file_path).stem}_p{page_num}_img{img_index}.{ext}"
                    img_path = os.path.join(self.image_output_dir, img_filename)

                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    # 创建图片描述块
                    desc = f"[图片: {img_filename}]\n"
                    desc += f"位置: 第{page_num}页\n"
                    desc += f"尺寸: {base_image.get('width', '?')}x{base_image.get('height', '?')}\n"
                    desc += f"颜色空间: {base_image.get('colorspace', '?')}"

                    chunks.append(self._create_chunk(
                        desc, ContentType.IMAGE.value,
                        {"image_path": img_path, "image_size": len(image_bytes)},
                        page_num
                    ))
            except Exception:
                continue

        return chunks

    def process_word(self, file_path: str) -> ProcessedDocument:
        """处理 Word 文件"""
        if Document is None:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        chunks = []
        all_text = []

        metadata = {
            "title": doc.core_properties.title or "",
            "author": doc.core_properties.author or "",
            "created": str(doc.core_properties.created) if doc.core_properties.created else "",
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables)
        }

        # 处理段落
        for para in doc.paragraphs:
            text = self._clean_text(para.text)
            if not text:
                continue

            all_text.append(text)

            # 判断段落样式
            style_name = para.style.name.lower() if para.style else ""

            if 'heading' in style_name or para.style and '标题' in style_name:
                chunk_type = ContentType.HEADING.value
                self._current_section = text
            elif text.startswith(('•', '-', '*', '1.', '2.', '一、', '二、')):
                chunk_type = ContentType.LIST.value
            else:
                chunk_type = ContentType.TEXT.value

            # 分块处理
            if len(text) > self.splitter.chunk_size:
                sub_chunks = self.splitter.split_text(text)
                for sub_text in sub_chunks:
                    chunks.append(self._create_chunk(
                        sub_text, chunk_type,
                        {"style": style_name}
                    ))
            else:
                chunks.append(self._create_chunk(
                    text, chunk_type,
                    {"style": style_name}
                ))

        # 处理表格
        for table_idx, table in enumerate(doc.tables, 1):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)

            table_text = self._extract_table_text(table_data)
            if table_text:
                chunks.append(self._create_chunk(
                    table_text, ContentType.TABLE.value,
                    {"table_index": table_idx, "row_count": len(table_data)}
                ))
                all_text.append(table_text)

        return self._build_document(file_path, "docx", 1, chunks, metadata, all_text)

    def process_text(self, file_path: str) -> ProcessedDocument:
        """处理纯文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

        text = self._clean_text(text)
        chunks = []

        # 尝试检测标题（Markdown 风格）
        lines = text.split('\n')
        current_text = []

        for line in lines:
            stripped = line.strip()

            # Markdown 标题
            if stripped.startswith('#'):
                if current_text:
                    content = '\n'.join(current_text)
                    sub_chunks = self.splitter.split_text(content)
                    for sub in sub_chunks:
                        chunks.append(self._create_chunk(
                            sub, ContentType.TEXT.value, {}
                        ))
                    current_text = []

                chunks.append(self._create_chunk(
                    stripped, ContentType.HEADING.value, {}
                ))
                self._current_section = stripped.lstrip('#').strip()

            # 代码块
            elif stripped.startswith('```'):
                if current_text:
                    content = '\n'.join(current_text)
                    sub_chunks = self.splitter.split_text(content)
                    for sub in sub_chunks:
                        chunks.append(self._create_chunk(
                            sub, ContentType.TEXT.value, {}
                        ))
                    current_text = []

                # 简单处理代码块
                chunks.append(self._create_chunk(
                    stripped, ContentType.CODE.value, {}
                ))

            else:
                current_text.append(line)

        if current_text:
            content = '\n'.join(current_text)
            sub_chunks = self.splitter.split_text(content)
            for sub in sub_chunks:
                chunks.append(self._create_chunk(
                    sub, ContentType.TEXT.value, {}
                ))

        metadata = {"line_count": len(lines), "char_count": len(text)}

        return self._build_document(file_path, "txt", 1, chunks, metadata, [text])

    def process_markdown(self, file_path: str) -> ProcessedDocument:
        """处理 Markdown 文件"""
        if markdown is None:
            return self.process_text(file_path)

        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()

        # 转换为 HTML 以便更好地解析
        html = markdown.markdown(md_text)

        if BeautifulSoup:
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text(separator='\n')
        else:
            text = md_text

        text = self._clean_text(text)
        chunks = []
        lines = md_text.split('\n')

        current_text = []
        in_code_block = False

        for line in lines:
            stripped = line.strip()

            if stripped.startswith('```'):
                in_code_block = not in_code_block
                if not in_code_block and current_text:
                    code_content = '\n'.join(current_text)
                    chunks.append(self._create_chunk(
                        code_content, ContentType.CODE.value,
                        {"language": stripped.strip('`').strip() or "unknown"}
                    ))
                    current_text = []
                continue

            if in_code_block:
                current_text.append(line)
                continue

            if stripped.startswith('#'):
                if current_text:
                    content = '\n'.join(current_text)
                    sub_chunks = self.splitter.split_text(content)
                    for sub in sub_chunks:
                        chunks.append(self._create_chunk(
                            sub, ContentType.TEXT.value, {}
                        ))
                    current_text = []

                level = len(stripped) - len(stripped.lstrip('#'))
                title = stripped.lstrip('#').strip()
                self._current_section = title
                chunks.append(self._create_chunk(
                    title, ContentType.HEADING.value,
                    {"heading_level": level}
                ))
            else:
                current_text.append(line)

        if current_text and not in_code_block:
            content = '\n'.join(current_text)
            sub_chunks = self.splitter.split_text(content)
            for sub in sub_chunks:
                chunks.append(self._create_chunk(sub, ContentType.TEXT.value, {}))

        metadata = {"original_format": "markdown"}
        return self._build_document(file_path, "md", 1, chunks, metadata, [text])

    def process_html(self, file_path: str) -> ProcessedDocument:
        """处理 HTML 文件"""
        if BeautifulSoup is None:
            raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4")

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()

        soup = BeautifulSoup(html, 'html.parser')

        # 移除 script 和 style
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()

        chunks = []

        # 提取标题
        for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            text = h.get_text(strip=True)
            if text:
                self._current_section = text
                chunks.append(self._create_chunk(
                    text, ContentType.HEADING.value,
                    {"tag": h.name}
                ))

        # 提取段落
        for p in soup.find_all(['p', 'div', 'article', 'section']):
            text = p.get_text(separator=' ', strip=True)
            if text and len(text) > 5:
                sub_chunks = self.splitter.split_text(text)
                for sub in sub_chunks:
                    chunks.append(self._create_chunk(
                        sub, ContentType.TEXT.value,
                        {"tag": p.name}
                    ))

        # 提取表格
        for table in soup.find_all('table'):
            rows = []
            for tr in table.find_all('tr'):
                row = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
                rows.append(row)

            table_text = self._extract_table_text(rows)
            if table_text:
                chunks.append(self._create_chunk(
                    table_text, ContentType.TABLE.value,
                    {"row_count": len(rows)}
                ))

        metadata = {"title": soup.title.string if soup.title else ""}
        all_text = [soup.get_text(separator='\n', strip=True)]

        return self._build_document(file_path, "html", 1, chunks, metadata, all_text)

    def process_ppt(self, file_path: str) -> ProcessedDocument:
        """处理 PPT 文件"""
        if Presentation is None:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        chunks = []
        all_text = []

        metadata = {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = self._clean_text(shape.text)
                    slide_texts.append(text)
                    all_text.append(text)

                # 提取表格
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows.append(row_data)

                    table_text = self._extract_table_text(rows)
                    if table_text:
                        chunks.append(self._create_chunk(
                            table_text, ContentType.TABLE.value,
                            {"slide": slide_num}, slide_num
                        ))

            if slide_texts:
                combined = "\n".join(slide_texts)
                sub_chunks = self.splitter.split_text(combined)
                for sub in sub_chunks:
                    chunks.append(self._create_chunk(
                        sub, ContentType.TEXT.value,
                        {"slide": slide_num}, slide_num
                    ))

        return self._build_document(file_path, "pptx", len(prs.slides), chunks, metadata, all_text)

    def process_excel(self, file_path: str) -> ProcessedDocument:
        """处理 Excel 文件"""
        if openpyxl is None:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, data_only=True)
        chunks = []
        all_text = []

        metadata = {
            "sheet_count": len(wb.sheetnames),
            "sheet_names": wb.sheetnames
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            self._current_section = f"工作表: {sheet_name}"

            # 提取所有单元格数据
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):
                    rows.append(row_data)

            if rows:
                table_text = self._extract_table_text(rows)
                sub_chunks = self.splitter.split_text(table_text)
                for sub in sub_chunks:
                    chunks.append(self._create_chunk(
                        sub, ContentType.TABLE.value,
                        {"sheet": sheet_name, "row_count": len(rows)}
                    ))
                all_text.append(table_text)

        wb.close()

        return self._build_document(file_path, "xlsx", 1, chunks, metadata, all_text)

    def _build_document(
            self,
            file_path: str,
            file_type: str,
            total_pages: int,
            chunks: List[Chunk],
            metadata: Dict[str, Any],
            all_text: List[str]
    ) -> ProcessedDocument:
        """构建最终文档对象"""
        full_text = " ".join(all_text)

        # 生成摘要统计
        word_count = len(full_text.split())
        char_count = len(full_text)

        # 提取关键词（简单实现）
        words = re.findall(r'\b[\u4e00-\u9fff]{2,4}\b|[a-zA-Z]{4,}', full_text)
        from collections import Counter
        top_words = Counter(words).most_common(20)

        summary = {
            "total_words": word_count,
            "total_chars": char_count,
            "chunk_types": dict(Counter(c.type for c in chunks)),
            "top_keywords": [w[0] for w in top_words[:10]],
            "avg_chunk_length": sum(len(c.content) for c in chunks) / len(chunks) if chunks else 0
        }

        return ProcessedDocument(
            source_file=os.path.basename(file_path),
            file_type=file_type,
            total_pages=total_pages,
            total_chunks=len(chunks),
            metadata=metadata,
            chunks=chunks,
            summary=summary
        )

    def process(self, file_path: str) -> ProcessedDocument:
        """根据文件类型自动选择处理方法"""
        path = Path(file_path)
        suffix = path.suffix.lower()

        processors = {
            '.pdf': self.process_pdf,
            '.docx': self.process_word,
            '.txt': self.process_text,
            '.md': self.process_markdown,
            '.markdown': self.process_markdown,
            '.html': self.process_html,
            '.htm': self.process_html,
            '.pptx': self.process_ppt,
            '.xlsx': self.process_excel,
            '.xls': self.process_excel,
        }

        if suffix in processors:
            return processors[suffix](file_path)
        else:
            # 尝试作为文本处理
            try:
                return self.process_text(file_path)
            except Exception as e:
                raise ValueError(f"不支持的文件类型: {suffix}，错误: {e}")


class AgentPromptBuilder:
    """构建 AI Agent 提示"""

    @staticmethod
    def build_qa_prompt(document: ProcessedDocument, question: str) -> str:
        """构建问答提示"""
        return f"""你是一个专业的文档分析助手。请基于以下文档内容回答问题。

## 文档信息
- 来源: {document.source_file}
- 类型: {document.file_type}
- 总块数: {document.total_chunks}

## 相关上下文
{document.to_agent_prompt(max_chunks=20)}

## 用户问题
{question}

请基于文档内容回答，如果文档中没有相关信息，请明确说明。"""

    @staticmethod
    def build_summary_prompt(document: ProcessedDocument) -> str:
        """构建摘要提示"""
        return f"""请对以下文档进行详细摘要。文档信息如下：

- 来源: {document.source_file}
- 总字数: {document.summary.get('total_words', 0)}
- 主要关键词: {', '.join(document.summary.get('top_keywords', [])[:5])}

## 文档内容
{document.to_agent_prompt(max_chunks=30)}

请提供：
1. 文档核心主题
2. 主要论点/内容概述
3. 关键数据或结论
4. 文档结构分析"""

    @staticmethod
    def build_extraction_prompt(document: ProcessedDocument, entity_type: str) -> str:
        """构建实体提取提示"""
        return f"""请从以下文档中提取所有{entity_type}。

## 文档内容
{document.to_agent_prompt(max_chunks=50)}

请以 JSON 格式输出提取结果，格式如下：
{{
    "entities": [
        {{"name": "实体名称", "context": "出现的上下文", "page": "页码"}}
    ]
}}"""


# ============ 使用示例 ============

def demo():
    """演示用法"""

    # 初始化处理器
    processor = DocumentProcessor(
        chunk_size=1000,  # 每块最大字符数
        chunk_overlap=150,  # 块间重叠字符数
        extract_images=False  # 是否提取图片
    )

    # 处理文件示例
    # result = processor.process("example.pdf")
    # result = processor.process("example.docx")
    # result = processor.process("example.txt")

    # 构建模拟数据用于演示
    demo_doc = ProcessedDocument(
        source_file="demo.pdf",
        file_type="pdf",
        total_pages=10,
        total_chunks=3,
        metadata={"author": "Demo", "title": "示例文档"},
        chunks=[
            Chunk(
                id="chunk_001",
                type="heading",
                content="第一章：人工智能概述",
                metadata={},
                page_num=1,
                section_title="第一章：人工智能概述",
                word_count=5,
                char_count=10
            ),
            Chunk(
                id="chunk_002",
                type="text",
                content="人工智能（AI）是计算机科学的一个分支，致力于创造能够模拟人类智能的系统。",
                metadata={},
                page_num=1,
                section_title="第一章：人工智能概述",
                word_count=20,
                char_count=40
            ),
            Chunk(
                id="chunk_003",
                type="table",
                content="模型 | 准确率 | 速度\nGPT-4 | 95% | 慢\nGPT-3.5 | 90% | 快",
                metadata={"table_index": 1},
                page_num=2,
                section_title="第二章：模型对比",
                word_count=10,
                char_count=50
            )
        ],
        summary={
            "total_words": 35,
            "total_chars": 100,
            "chunk_types": {"heading": 1, "text": 1, "table": 1},
            "top_keywords": ["人工智能", "模型", "GPT"],
            "avg_chunk_length": 33
        }
    )

    print("=" * 60)
    print("1. 输出为 JSON 格式")
    print("=" * 60)
    print(demo_doc.to_json(indent=2))

    print("\n" + "=" * 60)
    print("2. 输出为 Agent 提示格式")
    print("=" * 60)
    print(demo_doc.to_agent_prompt())

    print("\n" + "=" * 60)
    print("3. 构建问答提示")
    print("=" * 60)
    qa_prompt = AgentPromptBuilder.build_qa_prompt(demo_doc, "GPT-4 的准确率是多少？")
    print(qa_prompt)

    print("\n" + "=" * 60)
    print("4. 构建摘要提示")
    print("=" * 60)
    summary_prompt = AgentPromptBuilder.build_summary_prompt(demo_doc)
    print(summary_prompt)


if __name__ == "__main__":
    demo()